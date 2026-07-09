"""增量同步测试。"""

import tempfile
from pathlib import Path

import pytest

from app.documents.sync import DocumentSync, FileChange
from app.core.models import DocumentChunk, ChunkMetadata


def test_file_hash_calculation():
    """SHA-256 hash 计算正确。"""

    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.txt"
        test_file.write_text("hello", encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        hash1 = sync._file_hash(test_file)

        # 修改文件内容
        test_file.write_text("world", encoding="utf-8")
        hash2 = sync._file_hash(test_file)

        # hash 不同
        assert hash1 != hash2


def test_detect_changes_new_file():
    """检测新增文件变更。"""

    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.txt"
        test_file.write_text("content", encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        changes = sync.detect_changes()

        # 新文件 → add 变更
        assert len(changes) == 1
        assert changes[0].change_type == "add"
        assert changes[0].filename == "test.txt"


def test_sync_file_creates_chunks():
    """sync_file 创建 chunks 并写入活跃索引。"""

    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.txt"
        test_file.write_text("第一段内容\n\n第二段内容\n\n第三段内容", encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        count = sync.sync_file("test.txt")

        assert count >= 1
        assert "test.txt" in sync.active_chunks
        assert len(sync.active_chunks["test.txt"]) == count


def test_sync_file_markdown():
    """sync_file 对 Markdown 文件使用 chunk_markdown。"""

    with tempfile.TemporaryDirectory() as tmpdir:
        md_file = Path(tmpdir) / "药品说明书.md"
        md_file.write_text("# 标题1\n内容1\n\n## 标题2\n内容2", encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        count = sync.sync_file("药品说明书.md")

        assert count >= 1
        # Markdown chunks 应有 heading metadata
        chunks = sync.active_chunks["药品说明书.md"]
        assert any(c.metadata.heading for c in chunks)


def test_swap_moves_staging_to_active():
    """swap 原子性移动 staging → active。"""

    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.txt"
        test_file.write_text("内容", encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        # 手动添加到 staging
        sync.staging_chunks["manual.txt"] = [
            DocumentChunk(id="manual:1", source="manual.txt", content="手动内容")
        ]

        sync.swap()
        assert "manual.txt" in sync.active_chunks
        assert len(sync.staging_chunks) == 0


def test_get_total_chunk_count():
    """获取活跃索引总 chunk 数。"""

    sync = DocumentSync(knowledge_dir=Path("/tmp/nonexistent"))
    sync.active_chunks = {
        "a.md": [DocumentChunk(id="a:1", source="a.md", content="1")],
        "b.md": [
            DocumentChunk(id="b:1", source="b.md", content="1"),
            DocumentChunk(id="b:2", source="b.md", content="2"),
        ],
    }
    assert sync.get_total_chunk_count() == 3


def test_remove_unsynced_file_does_not_touch_indexes():
    """Removing an uploaded-but-unsynced file should not call external indexes."""

    class FailingIndex:
        def delete_chunks(self, source):
            raise AssertionError("external index should not be touched")

    sync = DocumentSync(
        knowledge_dir=Path("/tmp/nonexistent"),
        milvus_store=FailingIndex(),
        keyword_store=FailingIndex(),
    )

    sync._remove_chunks("uploaded-only.txt")


def test_sync_file_persists_index_state():
    """sync_file should persist chunk count for the lightweight document list."""

    from app.documents.index_state import load_index_state

    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.txt"
        test_file.write_text("medical document content " * 20, encoding="utf-8")

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        count = sync.sync_file("test.txt")

        state = load_index_state(Path(tmpdir))
        assert state["test.txt"]["chunk_count"] == count


def test_remove_chunks_clears_index_state():
    """_remove_chunks should clear persisted index state."""

    from app.documents.index_state import load_index_state, set_index_state

    with tempfile.TemporaryDirectory() as tmpdir:
        knowledge_dir = Path(tmpdir)
        set_index_state(knowledge_dir, "test.txt", 2)

        sync = DocumentSync(knowledge_dir=knowledge_dir)
        sync._remove_chunks("test.txt")

        assert "test.txt" not in load_index_state(knowledge_dir)

def _write_sync_fixture(path: Path) -> None:
    repeated_text = "医疗文档内容用于同步测试，包含适应症、用药剂量和注意事项。" * 12
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        prefix = "# 测试文档\n\n" if suffix == ".md" else ""
        path.write_text(prefix + repeated_text, encoding="utf-8")
        return

    if suffix == ".csv":
        path.write_text(
            "字段,说明\n适应症," + repeated_text + "\n注意事项," + repeated_text,
            encoding="utf-8",
        )
        return

    if suffix == ".xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "知识库"
        sheet.append(["字段", "说明"])
        sheet.append(["适应症", repeated_text])
        sheet.append(["注意事项", repeated_text])
        workbook.save(path)
        return

    if suffix == ".docx":
        from docx import Document

        document = Document()
        document.add_heading("测试文档", level=1)
        document.add_paragraph(repeated_text)
        document.add_paragraph(repeated_text)
        document.save(path)
        return

    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "测试文档"
        slide.placeholders[1].text = repeated_text
        presentation.save(path)
        return

    raise AssertionError(f"Unsupported fixture suffix: {suffix}")


@pytest.mark.parametrize("filename", [
    "medical.txt",
    "medical.md",
    "medical.csv",
    "medical.xlsx",
    "medical.docx",
    "medical.pptx",
])
def test_sync_file_supported_text_and_office_formats(filename):
    """sync_file should process common uploaded formats, not just plain text."""

    with tempfile.TemporaryDirectory() as tmpdir:
        file_path = Path(tmpdir) / filename
        _write_sync_fixture(file_path)

        sync = DocumentSync(knowledge_dir=Path(tmpdir))
        count = sync.sync_file(filename)

        assert count > 0
        assert filename in sync.active_chunks
        assert sync.active_chunks[filename][0].content.strip()
