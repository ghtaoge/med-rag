"""PPT 加载器。python-pptx 逐幻灯片提取文本。"""

from __future__ import annotations

from pathlib import Path

from app.documents.loader.base import DocumentLoader


class PptLoader(DocumentLoader):
    """PPT(.pptx) 加载器。"""

    def load(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("PPT 加载需要 python-pptx，请执行: pip install python-pptx")

        prs = Presentation(str(file_path))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            if texts:
                parts.append(f"# 幻灯片 {i + 1}\n\n" + "\n".join(texts))

        return "\n\n".join(parts)
